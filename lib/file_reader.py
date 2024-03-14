import os
from io import BytesIO
import io
from binascii import Error
import zipfile
import magic
from PyPDF2 import PdfReader
from PyPDF2.errors import PdfReadError
from openpyxl import load_workbook
from pptx import Presentation
from base64 import b64decode
import xml.etree.ElementTree as ET
from config.settings import ResultSetting, DatabaseSetting
from model.base64convert import Base64File
from model.base64_item import Base64Item
from utils.utils import Util
from utils.ansiprint import AnsiPrint
from lib.save_manager import SaveManager



class FileReader:
    def __init__(self) -> None:
        self._files_directory = ResultSetting().get_folder_output("files")
        self._invalid_b64_chars = '!"#$%&\'()*,-.:;<>?@[\\]^_`{|}~ '
        self._magic_numbers = {                                                                                                                                                                                                                      
            b"\x89\x50\x4E\x47\x0D\x0A\x1A\x0A": "png",
            b"\xFF\xD8\xFF": "jpeg",
            b"\x47\x49\x46\x38\x37\x61": "GIF", 
            b"\x50\x4B\x03\x04":"zip",
            b"%PDF-":"pdf"

        }    
    
    def _get_format(self, content:bytearray)->str:
        for magic_file, file_type in self._magic_numbers.items():
                    if content.startswith(magic_file):
                        return file_type
        return None
    
    def _read_pdf(self, content:bytearray, max_pages=2)->str:
        text = ""
        pdf = io.BytesIO(content)
        reader = PdfReader(pdf)
        total_pages = len(reader.pages)
        pages = min(total_pages, max_pages)
        for page in reader.pages:
            text+=page.extract_text()

        return text
        
    def _read_word(self, zip_file, document_name)->str:
        text=""
        with zip_file.open(document_name) as xml:
            
            tree = ET.parse(xml)
            root = tree.getroot()
            for elem in root.iter("{http://schemas.openxmlformats.org/wordprocessingml/2006/main}t"):
                if elem.text:                                                                                                                                                                                                              
                    text += elem.text + "\n"
        return text
    
    def _read_excel(self, content:bytearray)->str:
        xlsx = BytesIO(content)
        wb = load_workbook(xlsx)
        ws = wb.active
        text = ''
        
        for row in ws.iter_rows(values_only=True):
            for cell in row:
                if cell is not None:
                    text += str(cell) + ','
            
            if len([v for v in row if v is not None])>0:
                text+='\n'
                 
        
        return text
    
    def _read_ppt(self, content:bytes)->str:
        text=""
        with BytesIO(content) as stream:
            presentation = Presentation(stream)
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, 'text'):
                        text+=shape.text+'\n'
        
        return text
    
    def _get_extension(self, mime_type:str)->str:
        fragments = mime_type.split('/')
        return fragments[1], fragments[0]
         
    def get_zip_files(self, content:bytearray)->str:
        text = ""
        with zipfile.ZipFile(io.BytesIO(content), 'r') as zip_file:
            text = '\n'.join(zip_file.namelist())
        
        return text
    def get_file_type(self, content_bytes:bytes):
        #content_bytes = b64decode(content)
        return self.get_office(content_bytes)
    
        
    def get_office(self, content:bytearray):
        file_type:str = self._get_format(content)
        document_type:str = None
        text = ""
        if file_type is not None:
            if file_type == 'zip':
                    with zipfile.ZipFile(io.BytesIO(content), 'r') as zip_file:
                        if 'word/document.xml' in zip_file.namelist():
                                document_type = 'docx'
                                text = self._read_word(zip_file, 'word/document.xml')
                        elif 'xl/workbook.xml' in zip_file.namelist():
                                document_type = 'xlsx'
                                text = self._read_excel(content)
                        elif 'ppt/presentation.xml' in zip_file.namelist():
                                document_type = 'pptx'
                                text = self._read_ppt(content)
        
        return document_type, text
    
    
    def _is_base64_valid(self, text:str)->bool:
        if text is None or text.isdigit() or text == '':
           return False, None
        else:
           for character in text:
               if character in self._invalid_b64_chars:
                   return False, None
               
        return True, None
    
    def _save_file(self, content,format:str)->str:
        if content:
           save = SaveManager()
        
           filepath = save.save(content, format if format is not None else "", None, self._files_directory)
           return filepath
        
        return None
    
    def _complete_padding(self, text)->Base64Item:
        current_length = len(text)%4
        new_text = None
        if current_length > 0:
            padding_needed = 4-current_length
            new_text = text+("="*padding_needed)
        
        return Base64Item(original_text=text, fixed=new_text)
    
    def _decode(self, text:str)->Base64Item:
        item = Base64Item ()
        try:
            item.original_text = text
            item.fixed = self._complete_padding(text).fixed
            item.content = b64decode(text)
            
        except Error as be:
            if str(be) == 'Incorrect padding':
                fixed_text = self._complete_padding(text).fixed
                item.content = b64decode(fixed_text)
                
        except Exception as ex:
            AnsiPrint.print_debug(ex)
        
        return item
            

    def get_from_base64(self, text:str) -> Base64File:
        file = Base64File()
        try:
            if text:
                # if text.startswith("/9j/"):
                #     #text = self._complete_padding(text)
                #     print("")
                if text.startswith("data:image/"):
                    text = text.split(",")[1]

                content = None
                content_bytes = None
                
                valid, _ = self._is_base64_valid(text)
                
                if valid:
                    decoded = self._decode(text)
                    content_bytes = decoded.content
                    if content_bytes:
                        try:
                            
                            text_type = magic.from_buffer(content_bytes, mime=True)
                            ext, format = self._get_extension(text_type)
                            if ext == 'plain':
                                ext = 'txt'
                                if text[-2:] == '==':
                                    content = Util.get_readable_content(content_bytes)
                                elif Util.is_readable(content_bytes):
                                    content = content_bytes.decode('latin')
                                else:
                                    return file
                            elif ext == 'csv':
                                content = Util.get_readable_content(content_bytes)
                            elif ext == 'xml':
                                content = Util.get_readable_content(content_bytes)
                            elif ext == 'pdf':
                                content = self._read_pdf(content_bytes)
                            elif format == 'image': # We cannot convert to readable format
                                content = text
                            elif ext == 'zip':
                                content = self.get_zip_files(content_bytes)
                            elif format == 'application' and 'office' not in ext:
                                return file
                            
                        except UnicodeDecodeError as ue:
                            return file
                        
                        if content is None:
                            ext, content = self.get_file_type(content_bytes)
                        

                        file.extension = ext
                        file.content = content
                        file.filename = self._save_file(content_bytes, ext)
        except Exception as e:
            AnsiPrint.print_debug(e)
        
        return file
